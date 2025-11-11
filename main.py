from googleapiclient.discovery import build
from google.oauth2 import service_account
from fastapi import FastAPI, HTTPException, File, UploadFile
from docx import Document
from pptx import Presentation
from typing import Optional
import io
from pydantic import BaseModel
from contextlib import asynccontextmanager

# --- Configuración ---
SERVICE_ACCOUNT_FILE = "credentials-python.json"  # tu JSON de service account
SCOPES = ["https://www.googleapis.com/auth/drive.readonly"]
folder_id = "1JVV3OVjabbHIVvJZSb338w6ZrieDJ3IS"  # carpeta raíz con asignaturas

service = None
class WeeklyContentRequest(BaseModel):
    asignatura: str
    id_teoria: Optional[str] = None
    id_practica: Optional[str] = None
    id_laboratorio: Optional[str] = None

# --- Evento de inicio para autenticar una vez ---
@asynccontextmanager
async def lifespan(app: FastAPI):
    # Startup
    global service
    try:
        creds = service_account.Credentials.from_service_account_file(
            SERVICE_ACCOUNT_FILE, scopes=SCOPES
        )
        service = build("drive", "v3", credentials=creds)
        print("Servicio de Google Drive autenticado exitosamente.")
    except Exception as e:
        service = None
        print(f"Error fatal durante la autenticación: {e}")

    yield

    # Shutdown
    pass

app = FastAPI(lifespan=lifespan)


# --- Endpoint básico para health check ---
@app.get("/")
async def root():
    return {"message": "API de extracción de Google Drive funcionando correctamente"}


@app.get("/health")
async def health_check():
    return {"status": "ok", "service": "google_drive_extractor"}


# --- Modelo de Pydantic para los datos de entrada ---
class DriveQuery(BaseModel):
    ciclo_num: int
    semana_num: int


# Nuevo endpoint para recibir el binario y devolver el texto
@app.post("/convertir-documento")
async def convertir_documento_a_texto(file: UploadFile = File(...)):
    # El nombre del archivo se usa para determinar el tipo
    filename = file.filename
    file_extension = filename.split('.')[-1].lower()

    try:
        # 1. Leer el contenido binario usando el método estándar de FastAPI
        contents = await file.read()

        if file_extension == 'docx':
            # --- LÓGICA DOCX (Documentos) ---
            document = Document(io.BytesIO(contents))
            text_content = [paragraph.text for paragraph in document.paragraphs]

        elif file_extension == 'pptx':
            # --- LÓGICA PPTX (Presentaciones) ---
            prs = Presentation(io.BytesIO(contents))
            text_content = []

            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        # Aquí, debes asegurarte de obtener el texto real de los objetos
                        if shape.has_text_frame:
                            text_content.append(shape.text_frame.text)

            # Elimina entradas vacías de texto
            text_content = [t.strip() for t in text_content if t.strip()]

        else:
            # Si el archivo no es ni DOCX ni PPTX
            raise HTTPException(status_code=400,
                                detail=f"Formato no soportado: .{file_extension}. Solo se aceptan DOCX y PPTX.")

        return {"text": "\n".join(text_content)}

    except Exception as e:
        # Esto captura errores de procesamiento de archivos
        raise HTTPException(status_code=500,
                            detail=f"Error interno al procesar el archivo: {str(e)}. Asegura que 'python-docx' y 'python-pptx' estén instalados.")


# --- Endpoint POST principal ---
@app.post("/extraer-drive")
async def ejecutar_extraccion(query: DriveQuery):
    if not service:
        raise HTTPException(status_code=500,
                            detail="El servicio de Google Drive no está autenticado. Revisa los logs del contenedor.")

    try:
        # Llama a tu función principal con los datos del POST
        datos_completos = get_datos_ciclo(
            ciclo_num=query.ciclo_num,
            semana_num=query.semana_num
        )
        return datos_completos
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Ocurrió un error durante la extracción: {str(e)}")

@app.post("/extraer-contenido-semanal")
async def extraer_contenido_semanal(req: WeeklyContentRequest):
    if not service:
        raise HTTPException(status_code=500, detail="El servicio de Google Drive no está autenticado.")

    try:
        files = []
        seen = set()

        for folder_id_param in (req.id_teoria, req.id_practica, req.id_laboratorio):
            # skip None or empty strings
            if folder_id_param and folder_id_param.strip():
                archivos = get_files_in_folder(folder_id_param) or []
                for f in archivos:
                    fid = f.get("id")
                    if fid and fid not in seen:
                        seen.add(fid)
                        row = {
                        "id": fid,
                        }
                        files.append(row)

        return files
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error extrayendo contenido semanal: {e}")
# --- Funciones Auxiliares ---

# --- Nueva Función Auxiliar para la Evaluación Semanal Detallada (CORREGIDA) ---
def evaluate_weekly_folders(material_folder_id, semana_num):
    """
    Busca la carpeta 'Semana X' dentro de Teoria, Practica, y Laboratorio
    y devuelve el estado detallado, incluyendo el ID de la carpeta.
    """
    if not service: return None

    # Nota: He normalizado los nombres de las subcarpetas a minúsculas
    subfolder_types = ["teoría", "práctica", "laboratorio"]
    weekly_data = {}
    is_complete_week = False

    for subfolder_name in subfolder_types:
        # 1. Busca la carpeta de tipo (Teoría, Práctica, etc.)
        type_folder = find_item_in_folder(material_folder_id, subfolder_name,
                                          mime_type="application/vnd.google-apps.folder")

        # Inicializa el estado para el tipo actual
        weekly_data[subfolder_name] = {"estado": "No existe", "folder_id": None, "archivos": []}

        if type_folder:
            # 2. Busca la carpeta de la semana específica (ej: "Semana 6")
            semana_folder_name = f"semana {semana_num}"
            semana_folder = find_item_in_folder(type_folder['id'], semana_folder_name,
                                                mime_type="application/vnd.google-apps.folder")

            if semana_folder:
                # 3. Cuenta los archivos dentro de la carpeta semanal
                archivos_semana = get_files_in_folder(semana_folder['id'])

                estado = "Incompleto"
                if archivos_semana:
                    estado = "Completo"
                    # Si al menos un tipo tiene archivos, la semana general se considera completa para efectos de auditoría
                    is_complete_week = True

                weekly_data[subfolder_name] = {
                    "estado": estado,
                    "folder_id": semana_folder['id'],  # <-- ID de la carpeta semanal
                    "archivos": [{"nombre": f['name'], "id": f['id']} for f in archivos_semana]
                }
            else:
                # Si encuentra la carpeta Teoría/Práctica, pero no la Semana X
                weekly_data[subfolder_name] = {
                    "estado": "Falta Subir",
                    "folder_id": type_folder['id'], # Devuelve el ID de la carpeta padre si no encuentra la semana
                    "archivos": []
                }
            # Si no encuentra la carpeta de tipo (ej: no existe 'Laboratorio'), se queda con el estado inicial "No existe".

    return {"data": weekly_data, "overall_status": ("Completo" if is_complete_week else "Incompleto")}


def find_item_in_folder(parent_id, name_to_find, is_prefix=False, mime_type=None):
    """
    Busca un ítem de forma insensible a mayúsculas/minúsculas.
    Obtiene todos los ítems de la carpeta y los compara en Python para mayor precisión.
    """
    if not service: return None
    try:
        query = f"'{parent_id}' in parents and trashed=false"
        if mime_type:
            query += f" and mimeType = '{mime_type}'"

        results = service.files().list(q=query, fields="files(id, name, mimeType)").execute()
        items = results.get("files", [])

        # Normalizamos el nombre a buscar a minúsculas
        name_to_find_lower = name_to_find.lower()

        for item in items:
            # Comparamos todo en minúsculas y sin espacios extra
            item_name_lower = item['name'].lower().strip()

            if is_prefix:
                if item_name_lower.startswith(name_to_find_lower):
                    return item
            else:
                if item_name_lower == name_to_find_lower:
                    return item
        return None  # No se encontró el ítem
    except Exception as e:
        print(f"Error buscando '{name_to_find}' en '{parent_id}': {e}")
        return None


def get_all_folders_in_folder(parent_id):
    """Obtiene TODAS las subcarpetas de una carpeta padre."""
    if not service: return []
    try:
        query = f"'{parent_id}' in parents and trashed=false and mimeType = 'application/vnd.google-apps.folder'"
        results = service.files().list(q=query, fields="files(id, name)").execute()
        return results.get("files", [])
    except Exception as e:
        print(f"Error obteniendo carpetas de '{parent_id}': {e}")
        return []


def get_files_in_folder(folder_id):
    """Obtiene todos los archivos (no carpetas) de una carpeta."""
    if not service: return []
    try:
        query = f"'{folder_id}' in parents and trashed=false and mimeType != 'application/vnd.google-apps.folder'"
        results = service.files().list(q=query, fields="files(id, name)").execute()
        return results.get("files", [])
    except Exception as e:
        print(f"Error obteniendo archivos de la carpeta '{folder_id}': {e}")
        return []


# --- Función Principal Actualizada ---

def get_datos_ciclo(ciclo_num, semana_num):
    """
    Obtiene los datos de todas las asignaturas de un ciclo para una semana específica.
    Devuelve una lista con formato simplificado.
    """
    ciclo_folder = find_item_in_folder(folder_id, f"{ciclo_num} ", is_prefix=True,
                                       mime_type="application/vnd.google-apps.folder")
    if not ciclo_folder:
        return {"error": f"Ciclo {ciclo_num} no encontrado."}

    asignaturas = get_all_folders_in_folder(ciclo_folder['id'])
    if not asignaturas:
        return {"error": f"No se encontraron asignaturas en el ciclo {ciclo_num}."}

    resultados_ciclo = []
    for asignatura_folder in asignaturas:
        print(f"Procesando asignatura: {asignatura_folder['name']}...")

        # Inicializar IDs como None
        id_silabo = None
        id_teoria = None
        id_practica = None
        id_laboratorio = None

        # Búsqueda de sílabos
        silabo_folder = find_item_in_folder(asignatura_folder['id'], "1. silabo del curso",
                                            mime_type="application/vnd.google-apps.folder")
        if silabo_folder:
            id_silabo = silabo_folder['id']

        # Búsqueda de material de enseñanza
        material_folder = find_item_in_folder(asignatura_folder['id'], "2. material de enseñanza",
                                              mime_type="application/vnd.google-apps.folder")
        if material_folder:
            # Obtener los IDs de las carpetas de cada tipo para la semana específica
            semana_detalle = evaluate_weekly_folders(material_folder['id'], semana_num)

            if semana_detalle and "data" in semana_detalle:
                detalle_tipos = semana_detalle["data"]

                # Extraer los IDs de cada tipo
                if "teoría" in detalle_tipos and detalle_tipos["teoría"]["folder_id"]:
                    id_teoria = detalle_tipos["teoría"]["folder_id"]

                if "práctica" in detalle_tipos and detalle_tipos["práctica"]["folder_id"]:
                    id_practica = detalle_tipos["práctica"]["folder_id"]

                if "laboratorio" in detalle_tipos and detalle_tipos["laboratorio"]["folder_id"]:
                    id_laboratorio = detalle_tipos["laboratorio"]["folder_id"]

        # Crear el item con el formato solicitado
        item = {
            "asignatura": asignatura_folder['name'],
            "ciclo": ciclo_num,
            "semana": semana_num,
            "id_silabo": id_silabo,
            "id_teoria": id_teoria,
            "id_practica": id_practica,
            "id_laboratorio": id_laboratorio,
            "estado": "pendiente"
        }

        resultados_ciclo.append(item)

    return resultados_ciclo


if __name__ == "__main__":
    import uvicorn

    print("Iniciando servidor FastAPI en puerto 8000")
    uvicorn.run(app, host="0.0.0.0", port=8000)
